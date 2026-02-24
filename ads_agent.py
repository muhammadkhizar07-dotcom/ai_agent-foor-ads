import streamlit as st
import pandas as pd
import easyocr
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import google.generativeai as genai
import os
from dotenv import load_dotenv
import io

load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

model = genai.GenerativeModel('gemini-1.5-flash')  # For vision/text

st.title("Ads Performance AI Agent")

# Multiple CSV upload
csv_files = st.file_uploader("Upload CSVs (Meta/Google Ads)", type=["csv"], accept_multiple_files=True)

combined_df = None
if csv_files:
    all_data = []
    for uploaded_file in csv_files:
        try:
            df = pd.read_csv(uploaded_file, sep=',', engine='python', on_bad_lines="skip")
            df['source_file'] = uploaded_file.name
            all_data.append(df)
        except Exception as e:
            st.error(f"Error reading {uploaded_file.name}: {e}")
    if all_data:
        combined_df = pd.concat(all_data, ignore_index=True)
        st.subheader("Combined CSV Preview")
        st.dataframe(combined_df.head(10))

        # Dynamic metrics
        if "Impressions" in combined_df.columns and "Clicks" in combined_df.columns:
            combined_df["CTR (%)"] = (combined_df["Clicks"] / combined_df["Impressions"] * 100).fillna(0)
        if "Cost" in combined_df.columns and "Clicks" in combined_df.columns:
            combined_df["CPC"] = (combined_df["Cost"] / combined_df["Clicks"]).fillna(0)
        if "Cost" in combined_df.columns and "Conversions" in combined_df.columns:
            combined_df["CPA"] = (combined_df["Cost"] / combined_df["Conversions"]).fillna(0)

# Screenshot upload
image_file = st.file_uploader("Upload Screenshot", type=["png", "jpg", "jpeg"])
extracted_text = ""
if image_file:
    img = Image.open(image_file)
    st.image(img, caption="Screenshot")
    reader = easyocr.Reader(['en'])
    result = reader.readtext(img)
    extracted_text = " ".join([text for _, text, _ in result])
    st.write("Extracted Text:")
    st.text(extracted_text)

# User inputs (these will now show)
goal = st.text_input("Campaign Goal?")
budget = st.text_input("Budget?")
time_period = st.text_input("Time Period?")
changes = st.text_area("Changes Tested?")

# Generate button (always visible)
if st.button("Generate Outputs"):
    data_str = combined_df.to_string() if combined_df is not None else extracted_text
    if data_str:
        prompt = f"Analyze ads data: {data_str}. Goal: {goal}. Budget: {budget}. Period: {time_period}. Changes: {changes}. Write short case study on findings, actions, results. Suggest 3 LinkedIn post ideas for digital marketing skills. Generate 1 sample post. Suggest new strategies."
        
        response = model.generate_content(prompt)
        st.subheader("AI Outputs")
        st.write(response.text)

        # Stat visual
        if combined_df is not None and "CTR (%)" in combined_df.columns:
            plt.figure()
            combined_df["CTR (%)"].plot(kind='bar')
            plt.title("CTR by Entry")
            buf = io.BytesIO()
            plt.savefig(buf, format="png")
            buf.seek(0)
            st.image(buf)

        # PowerPoint
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Ads Case Study"
        content = slide.placeholders[1]
        content.text = response.text[:500]  # Snippet
        if 'buf' in locals():
            buf.seek(0)
            slide.shapes.add_picture(buf, Inches(1), Inches(2), Inches(8))
        prs.save("output.pptx")
        with open("output.pptx", "rb") as f:
            st.download_button("Download Slides", f, file_name="case_study.pptx")